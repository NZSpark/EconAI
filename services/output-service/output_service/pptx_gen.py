"""PPTX briefing generator (M7-20 through M7-24).

Generates presentation slides: cover, TOC, findings, recommendations, references.
"""

from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt


class PptxGenerator:
    """Generates briefing slides in .pptx format."""

    def __init__(self, institution_name: str | None = None) -> None:
        self._institution_name = institution_name or "PolicyAI 分析中心"

    def generate(
        self,
        title: str,
        sections: list[dict[str, Any]],
        citations: list[dict[str, Any]],
        metadata: dict[str, Any] | None = None,
    ) -> bytes:
        """生成 .pptx briefing.

        Args:
            title: Presentation title.
            sections: Analysis sections (used to extract findings).
            citations: Citation data for references slide.
            metadata: Optional metadata with subtitle, date, etc.

        Returns:
            PPTX file bytes.
        """
        meta = metadata or {}
        prs = Presentation()
        prs.slide_width = Cm(33.867)
        prs.slide_height = Cm(19.05)

        # Slide 1: Cover
        self._add_cover_slide(prs, title, meta)

        # Slide 2: TOC / Overview
        section_titles = [s.get("title", "") for s in sections if s.get("title")]
        self._add_toc_slide(prs, section_titles)

        # Slide 3-N: Key findings
        findings = meta.get("findings", [])
        if not findings:
            # Derive findings from sections
            for section in sections:
                content = section.get("content", "")
                if content:
                    findings.append(
                        {
                            "title": section.get("title", "关键发现"),
                            "points": [line.strip("- ") for line in content.split("\n") if line.strip()][:4],
                        }
                    )

        for finding in findings[:6]:
            self._add_finding_slide(prs, finding)

        # Recommendations / Conclusion slide
        recommendations = meta.get("recommendations", [])
        if recommendations:
            self._add_recommendations_slide(prs, recommendations)

        # Final slide: References
        self._add_references_slide(prs, citations)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def _add_cover_slide(self, prs: Any, title: str, meta: dict[str, Any]) -> None:
        """M7-20: Cover slide."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)

        # Set solid background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = title
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            slide.shapes.title.text_frame.paragraphs[0].font.bold = True
            slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Subtitle
        if len(slide.placeholders) > 1:
            ph = slide.placeholders[1]
            ph.text = meta.get("subtitle", "")
            ph.text_frame.paragraphs[0].font.size = Pt(18)
            ph.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            ph.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Date
        date_str = meta.get("date", "")
        if date_str:
            text_box = slide.shapes.add_textbox(Cm(4), Cm(15), Cm(25), Cm(2))
            tf = text_box.text_frame
            tf.paragraphs[0].text = date_str
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_toc_slide(self, prs: Any, section_titles: list[str]) -> None:
        """M7-21: Table of contents slide."""
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        if slide.shapes.title:
            slide.shapes.title.text = "目录"
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(24)
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
            slide.shapes.title.text_frame.paragraphs[0].font.bold = True

        if len(slide.placeholders) > 1:
            ph = slide.placeholders[1]
            tf = ph.text_frame
            tf.clear()
            for i, stitle in enumerate(section_titles[:8], 1):
                para = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
                para.text = f"{i}. {stitle}"
                para.font.size = Pt(16)
                para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                para.space_after = Pt(12)

    def _add_finding_slide(self, prs: Any, finding: dict[str, Any]) -> None:
        """M7-22: Key finding slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # Finding title
        if slide.shapes.title:
            slide.shapes.title.text = finding.get("title", "关键发现")
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(24)
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
            slide.shapes.title.text_frame.paragraphs[0].font.bold = True

        # Bullet points
        points = finding.get("points", [])
        if len(slide.placeholders) > 1 and points:
            ph = slide.placeholders[1]
            tf = ph.text_frame
            tf.clear()
            for i, point in enumerate(points):
                para = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                para.text = f"● {point}"
                para.font.size = Pt(16)
                para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                para.space_after = Pt(8)

        # Citation footer
        citation_text = finding.get("citation", "")
        if citation_text:
            text_box = slide.shapes.add_textbox(Cm(1.5), Cm(17.5), Cm(30), Cm(1))
            tf = text_box.text_frame
            tf.paragraphs[0].text = f"来源: {citation_text}"
            tf.paragraphs[0].font.size = Pt(9)
            tf.paragraphs[0].font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    def _add_recommendations_slide(self, prs: Any, recommendations: list[str]) -> None:
        """M7-23: Recommendations / Conclusion slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        if slide.shapes.title:
            slide.shapes.title.text = "政策建议"
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(24)
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
            slide.shapes.title.text_frame.paragraphs[0].font.bold = True

        if len(slide.placeholders) > 1:
            ph = slide.placeholders[1]
            tf = ph.text_frame
            tf.clear()
            for i, rec in enumerate(recommendations[:5], 1):
                para = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
                para.text = f"{i}. {rec}"
                para.font.size = Pt(16)
                para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                para.space_after = Pt(10)

    def _add_references_slide(self, prs: Any, citations: list[dict[str, Any]]) -> None:
        """M7-24: Full reference list slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        if slide.shapes.title:
            slide.shapes.title.text = "引用清单"
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(24)
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
            slide.shapes.title.text_frame.paragraphs[0].font.bold = True

        if len(slide.placeholders) > 1:
            ph = slide.placeholders[1]
            tf = ph.text_frame
            tf.clear()
            for i, cit in enumerate(citations[:20], 1):
                para = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
                ref_id = cit.get("ref_id", "")
                doc_title = cit.get("document_title", "")
                source_page = cit.get("source_page", cit.get("page_range", ""))
                para.text = f"[{i}] {doc_title}. {source_page} ({ref_id})"
                para.font.size = Pt(8)
                para.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
                para.space_after = Pt(4)
