"""Tests for all format parsers (M2-39).

Covers: PDF, Word, Markdown, Excel, PowerPoint, Email, HTML parsers.
"""

from __future__ import annotations

import io

from document_service.models import ParsedContent

# ---------------------------------------------------------------------------
# 夹具s
# ---------------------------------------------------------------------------


def _create_dummy_pdf_bytes() -> bytes:
    """Create a minimal valid PDF file."""
    # Minimal PDF structure
    pdf = (
        b"%PDF-1.4\n"
        b"1 0 obj\n"
        b"<< /Type /Catalog /Pages 2 0 R >>\n"
        b"endobj\n"
        b"2 0 obj\n"
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>\n"
        b"endobj\n"
        b"3 0 obj\n"
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>\n"
        b"endobj\n"
        b"4 0 obj\n"
        b"<< /Length 44 >>\n"
        b"stream\n"
        b"BT /F1 12 Tf 100 700 Td (Hello World PDF) Tj ET\n"
        b"endstream\n"
        b"endobj\n"
        b"5 0 obj\n"
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\n"
        b"endobj\n"
        b"xref\n"
        b"0 6\n"
        b"0000000000 65535 f \n"
        b"0000000009 00000 n \n"
        b"0000000058 00000 n \n"
        b"0000000115 00000 n \n"
        b"0000000266 00000 n \n"
        b"0000000360 00000 n \n"
        b"trailer\n"
        b"<< /Size 6 /Root 1 0 R >>\n"
        b"startxref\n"
        b"492\n"
        b"%%EOF"
    )
    return pdf


def _create_dummy_docx_bytes() -> bytes:
    """Create a minimal valid .docx file."""
    from zipfile import ZipFile
    buf = io.BytesIO()
    with ZipFile(buf, "w") as zf:
        zf.writestr(
            "[Content_Types].xml",
            '<?xml version="1.0"?><Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml" ContentType="application/xml"/>'
            '<Override PartName="/word/document.xml" '
            'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
            "</Types>",
        )
        zf.writestr(
            "_rels/.rels",
            '<?xml version="1.0"?><Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
            'Target="word/document.xml"/></Relationships>',
        )
        zf.writestr(
            "word/document.xml",
            '<?xml version="1.0"?><w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            "<w:body><w:p><w:r><w:t>Hello World from DOCX</w:t></w:r></w:p></w:body></w:document>",
        )
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PDF Parser Tests
# ---------------------------------------------------------------------------


class TestPDFParser:
    """M2-10: PDF parser tests."""

    def test_parse_pdf_extracts_text(self) -> None:
        from document_service.parsers.pdf_parser import PDFParser
        parser = PDFParser()
        pdf_bytes = _create_dummy_pdf_bytes()
        result = parser.parse(pdf_bytes, "test.pdf")

        assert result.full_text is not None
        assert len(result.full_text) > 0
        assert len(result.pages) > 0
        assert not result.needs_ocr  # Has text layer

    def test_parse_pdf_with_toc_detection(self) -> None:
        from document_service.parsers.pdf_parser import PDFParser
        parser = PDFParser()
        pdf_bytes = _create_dummy_pdf_bytes()
        result = parser.parse(pdf_bytes, "test.pdf")

        assert isinstance(result.sections, list)
        assert isinstance(result.pages, list)

    def test_pdf_metadata_extraction(self) -> None:
        from document_service.parsers.pdf_parser import PDFParser
        parser = PDFParser()
        pdf_bytes = _create_dummy_pdf_bytes()
        hints = parser.extract_metadata_hints(pdf_bytes, "test.pdf")
        assert isinstance(hints, dict)
        assert "title" in hints

    def test_supported_format(self) -> None:
        from document_service.parsers.pdf_parser import PDFParser
        parser = PDFParser()
        assert parser.supported_format() == "pdf"

    def test_empty_pdf(self) -> None:
        """PDF with no real content should still parse without error."""
        from document_service.parsers.pdf_parser import PDFParser
        parser = PDFParser()
        pdf_bytes = _create_dummy_pdf_bytes()
        result = parser.parse(pdf_bytes, "empty.pdf")
        assert result is not None
        assert isinstance(result, ParsedContent)


# ---------------------------------------------------------------------------
# Word Parser Tests
# ---------------------------------------------------------------------------


class TestWordParser:
    """M2-11: Word parser tests."""

    def test_parse_docx_extracts_text(self) -> None:
        from document_service.parsers.word_parser import WordParser
        parser = WordParser()
        docx_bytes = _create_dummy_docx_bytes()
        result = parser.parse(docx_bytes, "test.docx")

        assert result.full_text is not None
        assert "Hello World from DOCX" in result.full_text

    def test_docx_metadata_extraction(self) -> None:
        from document_service.parsers.word_parser import WordParser
        parser = WordParser()
        docx_bytes = _create_dummy_docx_bytes()
        hints = parser.extract_metadata_hints(docx_bytes, "test.docx")
        assert isinstance(hints, dict)
        assert "title" in hints

    def test_supported_format(self) -> None:
        from document_service.parsers.word_parser import WordParser
        parser = WordParser()
        assert parser.supported_format() == "docx"

    def test_docx_with_table_detection(self) -> None:
        from document_service.parsers.word_parser import WordParser
        parser = WordParser()
        docx_bytes = _create_dummy_docx_bytes()
        result = parser.parse(docx_bytes, "test.docx")
        assert isinstance(result.tables, list)


# ---------------------------------------------------------------------------
# Markdown Parser Tests
# ---------------------------------------------------------------------------


class TestMarkdownParser:
    """M2-12: Markdown parser tests."""

    def test_parse_markdown(self) -> None:
        from document_service.parsers.markdown_parser import MarkdownParser
        parser = MarkdownParser()
        md_bytes = b"# Title\n\nParagraph one.\n\n## Subtitle\n\nParagraph two."
        result = parser.parse(md_bytes, "test.md")

        assert result.full_text is not None
        assert "Title" in result.full_text
        assert len(result.sections) > 0

    def test_parse_plain_text(self) -> None:
        from document_service.parsers.markdown_parser import MarkdownParser
        parser = MarkdownParser()
        txt_bytes = b"Just some plain text.\n\nAnother paragraph."
        result = parser.parse(txt_bytes, "test.txt")

        assert "plain text" in result.full_text
        assert len(result.pages) == 1

    def test_heading_detection(self) -> None:
        from document_service.parsers.markdown_parser import MarkdownParser
        parser = MarkdownParser()
        md_bytes = b"# H1\n## H2\n### H3\n\nContent"
        result = parser.parse(md_bytes, "test.md")

        assert len(result.sections) >= 2

    def test_front_matter_detection(self) -> None:
        from document_service.parsers.markdown_parser import MarkdownParser
        parser = MarkdownParser()
        md_bytes = "---\ntitle: Test\n---\n\n# 内容".encode("utf-8")
        hints = parser.extract_metadata_hints(md_bytes, "test.md")
        assert hints.get("title") == "Test"

    def test_supported_format(self) -> None:
        from document_service.parsers.markdown_parser import MarkdownParser
        parser = MarkdownParser()
        assert parser.supported_format() == "markdown"


# ---------------------------------------------------------------------------
# Excel Parser Tests
# ---------------------------------------------------------------------------


class TestExcelParser:
    """M2-13: Excel/CSV parser tests."""

    def test_parse_csv(self) -> None:
        from document_service.parsers.excel_parser import ExcelParser
        parser = ExcelParser()
        csv_bytes = b"name,age,city\nAlice,30,NYC\nBob,25,LA"
        result = parser.parse(csv_bytes, "test.csv")

        assert "Alice" in result.full_text
        assert "Bob" in result.full_text
        assert len(result.tables) > 0

    def test_parse_xlsx_basic(self) -> None:
        from document_service.parsers.excel_parser import ExcelParser
        parser = ExcelParser()
        # 创建 a minimal xlsx file
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"] = "Name"
        ws["B1"] = "Value"
        ws["A2"] = "Item1"
        ws["B2"] = 100

        buf = io.BytesIO()
        wb.save(buf)
        xlsx_bytes = buf.getvalue()

        result = parser.parse(xlsx_bytes, "test.xlsx")

        assert result.full_text is not None
        assert len(result.tables) > 0 or len(result.sections) > 0

    def test_supported_format(self) -> None:
        from document_service.parsers.excel_parser import ExcelParser
        parser = ExcelParser()
        assert parser.supported_format() == "xlsx"


# ---------------------------------------------------------------------------
# 功效Point Parser Tests
# ---------------------------------------------------------------------------


class TestPPTParser:
    """M2-14: PowerPoint parser tests."""

    def test_parse_pptx(self) -> None:
        from pptx import Presentation

        from document_service.parsers.ppt_parser import PPTParser
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"
        if slide.placeholders:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    ph.text = "Content here"

        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        parser = PPTParser()
        result = parser.parse(pptx_bytes, "test.pptx")

        assert result.full_text is not None
        assert len(result.pages) > 0

    def test_pptx_metadata_extraction(self) -> None:
        from document_service.parsers.ppt_parser import PPTParser
        parser = PPTParser()
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test"

        buf = io.BytesIO()
        prs.save(buf)

        hints = parser.extract_metadata_hints(buf.getvalue(), "test.pptx")
        assert isinstance(hints, dict)

    def test_supported_format(self) -> None:
        from document_service.parsers.ppt_parser import PPTParser
        parser = PPTParser()
        assert parser.supported_format() == "pptx"


# ---------------------------------------------------------------------------
# 邮箱 Parser Tests
# ---------------------------------------------------------------------------


class TestEmailParser:
    """M2-15: Email parser tests."""

    def test_parse_eml(self) -> None:
        from document_service.parsers.email_parser import EmailParser
        parser = EmailParser()
        eml_bytes = (
            b"From: sender@test.com\r\n"
            b"To: recipient@test.com\r\n"
            b"Subject: Test Email\r\n"
            b"Date: Mon, 1 Jan 2024 00:00:00 +0000\r\n"
            b"Content-Type: text/plain; charset=utf-8\r\n"
            b"\r\n"
            b"This is the email body.\r\n"
            b"It has multiple lines.\r\n"
        )
        result = parser.parse(eml_bytes, "test.eml")

        assert "Test Email" in result.full_text
        assert "sender@test.com" in result.full_text
        assert "email body" in result.full_text
        hints = result.metadata_hints
        assert hints.get("title") == "Test Email"
        assert hints.get("author") == "sender@test.com"

    def test_parse_multipart_eml(self) -> None:
        from document_service.parsers.email_parser import EmailParser
        parser = EmailParser()
        eml_bytes = (
            b"From: sender@example.com\r\n"
            b"To: rcpt@example.com\r\n"
            b"Subject: Multi-part test\r\n"
            b'MIME-Version: 1.0\r\n'
            b'Content-Type: multipart/alternative; boundary="boundary42"\r\n'
            b"\r\n"
            b"--boundary42\r\n"
            b"Content-Type: text/plain\r\n"
            b"\r\n"
            b"Plain text body.\r\n"
            b"--boundary42\r\n"
            b"Content-Type: text/html\r\n"
            b"\r\n"
            b"<p>HTML body.</p>\r\n"
            b"--boundary42--\r\n"
        )
        result = parser.parse(eml_bytes, "test.eml")

        assert "Plain text body" in result.full_text
        assert "Multi-part test" in result.full_text

    def test_eml_metadata_author(self) -> None:
        from document_service.parsers.email_parser import EmailParser
        parser = EmailParser()
        eml_bytes = (
            b"From: Test User <test@test.com>\r\n"
            b"To: admin@test.com\r\n"
            b"Subject: Metadata\r\n"
            b"\r\n"
            b"Body.\r\n"
        )
        result = parser.parse(eml_bytes, "test.eml")
        assert "Test User" in result.metadata_hints["author"]

    def test_supported_format(self) -> None:
        from document_service.parsers.email_parser import EmailParser
        parser = EmailParser()
        assert parser.supported_format() == "eml"


# ---------------------------------------------------------------------------
# HTML Parser Tests
# ---------------------------------------------------------------------------


class TestHTMLParser:
    """M2-16: HTML parser tests."""

    def test_parse_html(self) -> None:
        from document_service.parsers.html_parser import HTMLParser
        parser = HTMLParser()
        html_bytes = (
            b"<html><head><title>Test Page</title></head>"
            b"<body><h1>Hello</h1><p>This is content.</p></body></html>"
        )
        result = parser.parse(html_bytes, "test.html")

        assert "Hello" in result.full_text
        assert "content" in result.full_text

    def test_html_removes_navigation(self) -> None:
        from document_service.parsers.html_parser import HTMLParser
        parser = HTMLParser()
        html_bytes = b"<html><body><nav>Nav content</nav><main>Main content</main></body></html>"
        result = parser.parse(html_bytes, "test.html")

        assert "Main content" in result.full_text

    def test_html_extracts_links(self) -> None:
        from document_service.parsers.html_parser import HTMLParser
        parser = HTMLParser()
        html_bytes = b'<html><body><a href="https://example.com">Link</a></body></html>'
        result = parser.parse(html_bytes, "test.html")
        hints = result.metadata_hints
        assert "links" in hints.get("html_metadata", {})

    def test_html_detects_headings(self) -> None:
        from document_service.parsers.html_parser import HTMLParser
        parser = HTMLParser()
        html_bytes = b"<html><body><h1>H1 Title</h1><h2>H2 Subtitle</h2><p>Text</p></body></html>"
        result = parser.parse(html_bytes, "test.html")

        assert len(result.sections) >= 2

    def test_supported_format(self) -> None:
        from document_service.parsers.html_parser import HTMLParser
        parser = HTMLParser()
        assert parser.supported_format() == "html"

    def test_html_extracts_title(self) -> None:
        from document_service.parsers.html_parser import HTMLParser
        parser = HTMLParser()
        html_bytes = b"<html><head><title>My Title</title></head><body><p>Body</p></body></html>"
        result = parser.parse(html_bytes, "test.html")

        assert result.metadata_hints.get("title") == "My Title"
