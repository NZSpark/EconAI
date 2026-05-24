"""PowerPoint parser using python-pptx (M2-14).

Extracts slide-by-slide text, notes, and embedded images with OCR recognition.
"""

from __future__ import annotations

import io
import logging
from typing import Any

from document_service.models import PageContent, ParsedContent, SectionInfo
from document_service.parsers.base import BaseParser
from document_service.parsers.image_extractor import extract_images_from_pptx

logger = logging.getLogger(__name__)


class PPTParser(BaseParser):
    """Parse .pptx files using python-pptx."""

    def supported_format(self) -> str:
        return "pptx"

    def parse(self, file_data: bytes, filename: str) -> ParsedContent:
        try:
            from pptx import Presentation
        except ImportError as e:
            raise ImportError("python-pptx is required for PowerPoint parsing") from e

        prs = Presentation(io.BytesIO(file_data))
        pages: list[PageContent] = []
        all_text_parts: list[str] = []
        sections: list[SectionInfo] = []
        all_tables: list[dict[str, Any]] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text_parts: list[str] = []
            title = ""

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text_parts.append(text)
                            # First text frame is often the title
                            if not title:
                                title = text

                if shape.has_table:
                    table = shape.table
                    rows = [[cell.text for cell in row.cells] for row in table.rows]
                    all_tables.append({"slide": slide_num, "rows": rows})

            # Extract notes
            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text

            slide_text = "\n".join(slide_text_parts)
            if notes_text:
                slide_text += f"\n[Notes] {notes_text}"

            pages.append(PageContent(page_number=slide_num, text=slide_text))
            all_text_parts.append(f"--- Slide {slide_num} ---\n{slide_text}")

            if title:
                sections.append(SectionInfo(title=title, level=1, page_start=slide_num))

        full_text = "\n\n".join(all_text_parts)

        # ---- Extract and OCR embedded images per slide ----
        ocr_images = extract_images_from_pptx(file_data)
        image_text_by_slide: dict[int, list[str]] = {}
        for img in ocr_images:
            slide_num = img.get("page", 1)
            ocr_text = img.get("ocr_text", "")
            if ocr_text:
                image_text_by_slide.setdefault(slide_num, []).append(
                    f"[Image {img.get('image_index', 0)} OCR]: {ocr_text}"
                )

        # Merge image OCR into corresponding slide content
        for page_content in pages:
            pn = page_content.page_number
            if pn in image_text_by_slide:
                extra = "\n\n".join(image_text_by_slide[pn])
                page_content.text = (page_content.text + "\n\n" + extra).strip()

        if image_text_by_slide:
            # Rebuild full_text with enriched slide content
            enriched_parts = []
            for page_content in pages:
                enriched_parts.append(f"--- Slide {page_content.page_number} ---\n{page_content.text}")
            full_text = "\n\n".join(enriched_parts)
            logger.info("PPTX: OCR'd %d images across slides", len(ocr_images))
        # ---- End image extraction ----

        return ParsedContent(
            full_text=full_text,
            pages=pages,
            tables=all_tables,
            sections=sections,
            metadata_hints=self.extract_metadata_hints(file_data, filename),
            needs_ocr=False,
            ocr_images=ocr_images,
        )

    def extract_metadata_hints(self, file_data: bytes, filename: str) -> dict[str, Any]:
        """Extract PPT metadata."""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_data))
            props = prs.core_properties
            return {
                "title": props.title or "",
                "author": props.author or "",
                "date": str(props.created) if props.created else "",
                "source": props.last_modified_by or "",
                "ppt_metadata": {
                    "title": props.title,
                    "author": props.author,
                    "created": str(props.created) if props.created else None,
                    "modified": str(props.modified) if props.modified else None,
                    "last_modified_by": props.last_modified_by,
                    "category": props.category,
                    "subject": props.subject,
                },
            }
        except Exception:
            return {"title": filename}
