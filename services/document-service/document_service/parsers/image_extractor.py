"""Shared image extraction and OCR utility for multi-format document parsers.

Extracts embedded images from PDF, DOCX, PPTX, and HTML documents,
then runs Tesseract OCR (chi_sim+eng) on each image.

Supports:
  - PDF: Embedded images via PyMuPDF page.get_images() + extract_image()
  - DOCX: Embedded images via python-docx part relationships
  - PPTX: Slide images via python-pptx shape.image
  - HTML: Data-URI images (<img src="data:image/...">)

Design notes:
  - OCR is optional: gracefully degrades if pytesseract/tesseract not installed.
  - All extraction methods return [{"page", "image_index", "ocr_text", "format", "width", "height"}, ...].
  - Language default: chi_sim+eng (Chinese simplified + English).
"""

from __future__ import annotations

import base64
import io
import logging
import re
from typing import Any

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Shared OCR helper
# ---------------------------------------------------------------------------


def ocr_image_bytes(image_bytes: bytes, language: str = "chi_sim+eng") -> str:
    """Run Tesseract OCR on raw image bytes.

    Args:
        image_bytes: Raw image data (PNG, JPEG, etc.).
        language: Tesseract language string (default chi_sim+eng).

    Returns:
        Extracted text, or a descriptive message if OCR is unavailable/fails.
    """
    try:
        import pytesseract
        from PIL import Image

        img = Image.open(io.BytesIO(image_bytes))
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")

        text: str = pytesseract.image_to_string(
            img,
            lang=language,
            config="--psm 6",  # Assume uniform block of text
        )
        return text.strip()
    except ImportError:
        logger.warning("pytesseract not available; returning empty OCR result for image")
        return "[OCR not available: pytesseract not installed]"
    except Exception as e:
        logger.warning("Tesseract OCR error: %s", e)
        return f"[OCR Error: {e}]"


# ---------------------------------------------------------------------------
# PDF image extraction
# ---------------------------------------------------------------------------


def extract_images_from_pdf(
    file_data: bytes, language: str = "chi_sim+eng"
) -> list[dict[str, Any]]:
    """Extract and OCR all embedded images from a PDF document.

    Uses PyMuPDF (fitz) to find embedded images on each page via
    page.get_images(full=True), then runs OCR on each extracted image.

    Args:
        file_data: Raw PDF file bytes.
        language: Tesseract language string.

    Returns:
        List of dicts with keys: page, image_index, ocr_text, format, width, height.
    """
    results: list[dict[str, Any]] = []

    try:
        import fitz
    except ImportError:
        logger.warning("PyMuPDF (fitz) not available for PDF image extraction")
        return results

    try:
        doc = fitz.open(stream=file_data, filetype="pdf")
    except Exception as e:
        logger.warning("Cannot open PDF for image extraction: %s", e)
        return results

    try:
        for page_num in range(len(doc)):
            page = doc[page_num]
            try:
                image_list = page.get_images(full=True)
            except Exception:
                continue

            if not image_list:
                continue

            for img_idx, img_info in enumerate(image_list):
                try:
                    xref = img_info[0]  # xref number
                    base_image = doc.extract_image(xref)
                    img_bytes = base_image.get("image")
                    if not img_bytes:
                        continue

                    img_ext = base_image.get("ext", "png")
                    img_width = base_image.get("width", 0)
                    img_height = base_image.get("height", 0)

                    # Skip tiny images (likely icons, decorations)
                    if img_width < 20 or img_height < 20:
                        continue

                    ocr_text = ocr_image_bytes(img_bytes, language)

                    results.append({
                        "page": page_num + 1,
                        "image_index": img_idx,
                        "ocr_text": ocr_text,
                        "format": img_ext,
                        "width": img_width,
                        "height": img_height,
                    })
                except Exception as e:
                    logger.debug("Failed to extract/OCR image %d on page %d: %s", img_idx, page_num + 1, e)

    finally:
        doc.close()

    if results:
        logger.info("Extracted and OCR'd %d images from PDF", len(results))
    return results


# ---------------------------------------------------------------------------
# DOCX image extraction
# ---------------------------------------------------------------------------


def extract_images_from_docx(
    file_data: bytes, language: str = "chi_sim+eng"
) -> list[dict[str, Any]]:
    """Extract and OCR all embedded images from a DOCX document.

    Iterates part relationships in the python-docx document object
    to find image parts, then runs OCR on each.

    Args:
        file_data: Raw DOCX file bytes.
        language: Tesseract language string.

    Returns:
        List of dicts with keys: page, image_index, ocr_text, format, width, height.
    """
    results: list[dict[str, Any]] = []

    try:
        from docx import Document
    except ImportError:
        logger.warning("python-docx not available for DOCX image extraction")
        return results

    try:
        doc = Document(io.BytesIO(file_data))
    except Exception as e:
        logger.warning("Cannot open DOCX for image extraction: %s", e)
        return results

    try:
        img_idx = 0
        for rel in doc.part.rels.values():
            reltype = str(rel.reltype) if hasattr(rel, "reltype") else ""
            if "image" not in reltype.lower():
                continue

            try:
                image_part = rel.target_part
                img_bytes = image_part.blob
                content_type = getattr(image_part, "content_type", "image/png")
            except Exception as e:
                logger.debug("Skipping DOCX image relationship: %s", e)
                continue

            if not img_bytes:
                continue

            # Determine image format from content type
            fmt_map = {
                "image/png": "png", "image/jpeg": "jpeg",
                "image/gif": "gif", "image/bmp": "bmp",
                "image/tiff": "tiff", "image/webp": "webp",
            }
            img_format = fmt_map.get(content_type, "png")

            # Get image dimensions via PIL
            width, height = 0, 0
            try:
                from PIL import Image

                pil_img = Image.open(io.BytesIO(img_bytes))
                width, height = pil_img.size
            except Exception:
                pass

            # Skip tiny images
            if width < 20 or height < 20:
                continue

            ocr_text = ocr_image_bytes(img_bytes, language)

            results.append({
                "page": 1,
                "image_index": img_idx,
                "ocr_text": ocr_text,
                "format": img_format,
                "width": width,
                "height": height,
            })
            img_idx += 1

    except Exception as e:
        logger.warning("DOCX image extraction failed: %s", e)

    if results:
        logger.info("Extracted and OCR'd %d images from DOCX", len(results))
    return results


# ---------------------------------------------------------------------------
# PPTX image extraction
# ---------------------------------------------------------------------------


def extract_images_from_pptx(
    file_data: bytes, language: str = "chi_sim+eng"
) -> list[dict[str, Any]]:
    """Extract and OCR all embedded images from a PPTX presentation.

    Iterates each slide's shapes, finds PICTURE shapes, and runs OCR.

    Args:
        file_data: Raw PPTX file bytes.
        language: Tesseract language string.

    Returns:
        List of dicts with keys: page, image_index, ocr_text, format, width, height.
    """
    results: list[dict[str, Any]] = []

    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        logger.warning("python-pptx not available for PPTX image extraction")
        return results

    try:
        prs = Presentation(io.BytesIO(file_data))
    except Exception as e:
        logger.warning("Cannot open PPTX for image extraction: %s", e)
        return results

    try:
        img_idx = 0
        for slide_num, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue

                try:
                    image = shape.image
                    img_bytes = image.blob
                    content_type = getattr(image, "content_type", "image/png")
                except Exception as e:
                    logger.debug("Skipping PPTX image shape on slide %d: %s", slide_num, e)
                    continue

                if not img_bytes:
                    continue

                fmt_map = {
                    "image/png": "png", "image/jpeg": "jpeg",
                    "image/gif": "gif", "image/bmp": "bmp",
                }
                img_format = fmt_map.get(content_type, "png")

                width, height = 0, 0
                try:
                    from PIL import Image

                    pil_img = Image.open(io.BytesIO(img_bytes))
                    width, height = pil_img.size
                except Exception:
                    pass

                if width < 20 or height < 20:
                    continue

                ocr_text = ocr_image_bytes(img_bytes, language)

                results.append({
                    "page": slide_num,
                    "image_index": img_idx,
                    "ocr_text": ocr_text,
                    "format": img_format,
                    "width": width,
                    "height": height,
                })
                img_idx += 1

    except Exception as e:
        logger.warning("PPTX image extraction failed: %s", e)

    if results:
        logger.info("Extracted and OCR'd %d images from PPTX", len(results))
    return results


# ---------------------------------------------------------------------------
# HTML data-URI image extraction
# ---------------------------------------------------------------------------


_DATA_URI_PATTERN = re.compile(
    r'<img[^>]+src="data:image/(\w+);base64,([^"]+)"',
    re.IGNORECASE,
)


def extract_images_from_html(
    file_data: bytes, language: str = "chi_sim+eng"
) -> list[dict[str, Any]]:
    """Extract and OCR base64-encoded data-URI images from HTML content.

    Args:
        file_data: Raw HTML file bytes.
        language: Tesseract language string.

    Returns:
        List of dicts with keys: page, image_index, ocr_text, format, width, height.
    """
    results: list[dict[str, Any]] = []

    try:
        html_text = file_data.decode("utf-8")
    except UnicodeDecodeError:
        html_text = file_data.decode("latin-1", errors="replace")

    matches = _DATA_URI_PATTERN.findall(html_text)
    if not matches:
        return results

    for img_idx, (img_format, b64_data) in enumerate(matches):
        try:
            img_bytes = base64.b64decode(b64_data)
        except Exception as e:
            logger.debug("Failed to decode base64 image %d in HTML: %s", img_idx, e)
            continue

        if len(img_bytes) < 100:
            continue  # Too small to be meaningful

        ocr_text = ocr_image_bytes(img_bytes, language)

        results.append({
            "page": 1,
            "image_index": img_idx,
            "ocr_text": ocr_text,
            "format": img_format.lower(),
            "width": 0,
            "height": 0,
        })

    if results:
        logger.info("Extracted and OCR'd %d data-URI images from HTML", len(results))
    return results
